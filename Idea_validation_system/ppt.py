"""
ppt.py — Pitch Deck Generator for ThynxAI (Offline Edition)

Generates a professional pitch deck from the analysis dict using python-pptx.
No template file required — builds slides programmatically with a clean dark theme.
"""

import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─── Theme Colors ──────────────────────────────────────────────────────────────

BG_DARK    = RGBColor(0x0F, 0x0F, 0x1A)   # near-black background
BG_CARD    = RGBColor(0x1E, 0x1E, 0x2E)   # card background
ACCENT     = RGBColor(0x89, 0xB4, 0xFA)   # soft blue accent
ACCENT2    = RGBColor(0xA6, 0xE3, 0xA1)   # green for positives
TITLE_COL  = RGBColor(0xCD, 0xD6, 0xF4)   # off-white title
BODY_COL   = RGBColor(0xBA, 0xC2, 0xDE)   # muted body text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


# ─── Slide helpers ─────────────────────────────────────────────────────────────

def _set_bg(slide, prs, color: RGBColor):
    """Fill the entire slide background with a solid colour."""
    from pptx.oxml.ns import qn
    from lxml import etree
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height, font_size=14, bold=False, color=BODY_COL, align=PP_ALIGN.LEFT):
    """Helper to add a styled textbox."""
    if isinstance(text, list):
        text = ", ".join(str(x) for x in text)
    elif isinstance(text, dict):
        text = ", ".join(f"{k}: {v}" for k, v in text.items())
    elif text is None:
        text = "—"
    elif not isinstance(text, str):
        text = str(text)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = color
    return txBox


def _add_slide(prs: Presentation) -> object:
    """Add a blank slide with dark background."""
    layout = prs.slide_layouts[6]   # blank layout
    slide  = prs.slides.add_slide(layout)
    _set_bg(slide, prs, BG_DARK)
    return slide


def _header_bar(slide, prs, title: str):
    """Add a top accent bar + slide title."""
    W = prs.slide_width
    # Thin accent bar at top
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, W, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    _add_textbox(slide, title,
                 Inches(0.4), Inches(0.12),
                 W - Inches(0.8), Inches(0.55),
                 font_size=22, bold=True, color=TITLE_COL)


def _bullet_box(slide, items: list, left, top, width, height,
                bullet="•", color=BODY_COL, font_size=14):
    lines = "\n".join(f"{bullet}  {i}" for i in items) if items else "—"
    _add_textbox(slide, lines, left, top, width, height,
                 font_size=font_size, color=color)


# ─── Individual slide builders ─────────────────────────────────────────────────

def _slide_title(prs, analysis: dict):
    slide = _add_slide(prs)
    W, H  = prs.slide_width, prs.slide_height
    overall = analysis.get("overall", {})
    score   = overall.get("score", "?")

    # Big gradient-feel rectangle
    rect = slide.shapes.add_shape(1, 0, Inches(2.5), W, Inches(2.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_CARD
    rect.line.fill.background()

    _add_textbox(slide, "🚀 ThynxAI Idea Lab",
                 Inches(0.5), Inches(0.4), W - Inches(1), Inches(0.6),
                 font_size=14, color=ACCENT, bold=False)

    idea = analysis.get("idea_summary", "Startup Idea Validation")
    _add_textbox(slide, idea,
                 Inches(0.5), Inches(1.0), W - Inches(1), Inches(1.2),
                 font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_textbox(slide, f"Founder: {analysis.get('founder_name', '—')}",
                 Inches(0.5), Inches(2.8), Inches(4), Inches(0.4),
                 font_size=14, color=ACCENT2)

    _add_textbox(slide, f"Overall Score: {score}/10",
                 Inches(4.5), Inches(2.8), Inches(3), Inches(0.4),
                 font_size=14, bold=True, color=ACCENT)

    verdict = overall.get("final_verdict", "")
    _add_textbox(slide, verdict,
                 Inches(0.5), Inches(3.4), W - Inches(1), Inches(1.0),
                 font_size=13, color=BODY_COL, align=PP_ALIGN.CENTER)

    pitch = analysis.get("proposed_solution", {}).get("one_line_pitch", "")
    _add_textbox(slide, f'"{pitch}"',
                 Inches(0.5), Inches(4.5), W - Inches(1), Inches(0.7),
                 font_size=14, color=ACCENT2, align=PP_ALIGN.CENTER)


def _slide_problem(prs, analysis: dict):
    slide = _add_slide(prs)
    W     = prs.slide_width
    _header_bar(slide, prs, "1️⃣  Problem Statement")
    prob  = analysis.get("problem_statement", {})

    _add_textbox(slide, prob.get("description", "—"),
                 Inches(0.4), Inches(0.85), W - Inches(0.8), Inches(0.55),
                 font_size=14, color=WHITE)

    _add_textbox(slide, "Target Audience",
                 Inches(0.4), Inches(1.5), Inches(3.5), Inches(0.35),
                 font_size=12, bold=True, color=ACCENT)
    _add_textbox(slide, prob.get("target_audience", "—"),
                 Inches(0.4), Inches(1.85), Inches(3.5), Inches(0.4),
                 font_size=12, color=BODY_COL)

    _add_textbox(slide, "Market Size",
                 Inches(4.2), Inches(1.5), Inches(3.5), Inches(0.35),
                 font_size=12, bold=True, color=ACCENT)
    _add_textbox(slide, prob.get("market_size_hint", "—"),
                 Inches(4.2), Inches(1.85), Inches(3.5), Inches(0.4),
                 font_size=12, color=BODY_COL)

    _add_textbox(slide, "Pain Points",
                 Inches(0.4), Inches(2.4), W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _bullet_box(slide, prob.get("pain_points", []),
                Inches(0.4), Inches(2.75), W - Inches(0.8), Inches(1.5))

    _add_textbox(slide, "Real World Example",
                 Inches(0.4), Inches(4.3), W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _add_textbox(slide, prob.get("real_world_example", "—"),
                 Inches(0.4), Inches(4.65), W - Inches(0.8), Inches(0.7),
                 font_size=12, color=BODY_COL)


def _slide_solution(prs, analysis: dict):
    slide = _add_slide(prs)
    W     = prs.slide_width
    _header_bar(slide, prs, "2️⃣  Proposed Solution")
    sol   = analysis.get("proposed_solution", {})

    _add_textbox(slide, sol.get("simple_explanation", "—"),
                 Inches(0.4), Inches(0.85), W - Inches(0.8), Inches(0.65),
                 font_size=14, color=WHITE)

    _add_textbox(slide, "How It Works",
                 Inches(0.4), Inches(1.6), Inches(4), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _bullet_box(slide, sol.get("step_by_step_how_it_works", []),
                Inches(0.4), Inches(1.95), Inches(4), Inches(2.0),
                bullet="→")

    _add_textbox(slide, "Key Features",
                 Inches(4.4), Inches(1.6), Inches(3.2), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _bullet_box(slide, sol.get("key_features", []),
                Inches(4.4), Inches(1.95), Inches(3.2), Inches(2.0))

    _add_textbox(slide, "⚡ Unfair Advantage",
                 Inches(0.4), Inches(4.1), W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT2)
    _add_textbox(slide, sol.get("unfair_advantage", "—"),
                 Inches(0.4), Inches(4.45), W - Inches(0.8), Inches(0.7),
                 font_size=13, color=ACCENT2)


def _slide_market(prs, analysis: dict):
    slide = _add_slide(prs)
    W     = prs.slide_width
    _header_bar(slide, prs, "3️⃣  Market Landscape")
    mkt   = analysis.get("market_landscape", {})
    innov = analysis.get("core_innovation", {})

    _add_textbox(slide, "Competition Level",
                 Inches(0.4), Inches(0.85), Inches(3.5), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _add_textbox(slide, mkt.get("competition_level", "—"),
                 Inches(0.4), Inches(1.2), Inches(3.5), Inches(0.4),
                 font_size=16, bold=True, color=WHITE)

    _add_textbox(slide, "Innovation Type",
                 Inches(4.2), Inches(0.85), Inches(3.5), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _add_textbox(slide, innov.get("innovation_type", "—"),
                 Inches(4.2), Inches(1.2), Inches(3.5), Inches(0.4),
                 font_size=16, bold=True, color=WHITE)

    _add_textbox(slide, "Similar Solutions",
                 Inches(0.4), Inches(1.8), W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _add_textbox(slide, mkt.get("similar_solutions", "—"),
                 Inches(0.4), Inches(2.15), W - Inches(0.8), Inches(0.65),
                 font_size=12, color=BODY_COL)

    _add_textbox(slide, "Market Gap",
                 Inches(0.4), Inches(2.95), W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT2)
    _add_textbox(slide, mkt.get("market_gap", "—"),
                 Inches(0.4), Inches(3.3), W - Inches(0.8), Inches(0.75),
                 font_size=13, color=ACCENT2)

    _add_textbox(slide, "Core Uniqueness",
                 Inches(0.4), Inches(4.15), W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _add_textbox(slide, innov.get("uniqueness", "—"),
                 Inches(0.4), Inches(4.5), W - Inches(0.8), Inches(0.7),
                 font_size=12, color=BODY_COL)


def _slide_scores(prs, analysis: dict):
    slide  = _add_slide(prs)
    W      = prs.slide_width
    _header_bar(slide, prs, "4️⃣  Scores & Evaluation")
    scores = analysis.get("scores", {})
    labels = {
        "market_feasibility":  "Market Feasibility",
        "marketing_potential": "Marketing Potential",
        "scalability":         "Scalability",
        "revenue_potential":   "Revenue Potential",
        "technical_complexity":"Technical Complexity",
        "execution_risk":      "Execution Risk",
    }

    col_w = Inches(3.7)
    for idx, (key, label) in enumerate(labels.items()):
        col  = idx % 2
        row  = idx // 2
        x    = Inches(0.4) + col * (col_w + Inches(0.1))
        y    = Inches(0.9) + row * Inches(1.4)
        data = scores.get(key, {})
        score_val = data.get("score", "?")
        reason    = data.get("reasoning", "")[:90]

        _add_textbox(slide, f"{label}  —  {score_val}/10",
                     x, y, col_w, Inches(0.4),
                     font_size=13, bold=True, color=ACCENT)
        _add_textbox(slide, reason,
                     x, y + Inches(0.38), col_w, Inches(0.85),
                     font_size=11, color=BODY_COL)

    overall = analysis.get("overall", {})
    _add_textbox(slide, f"Overall: {overall.get('score','?')}/10",
                 Inches(0.4), Inches(5.15), W - Inches(0.8), Inches(0.4),
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def _slide_roadmap(prs, analysis: dict):
    slide   = _add_slide(prs)
    W       = prs.slide_width
    _header_bar(slide, prs, "5️⃣  Roadmap & Support")
    overall = analysis.get("overall", {})
    support = analysis.get("support_required", {})
    tech    = analysis.get("tech_stack", {})

    readiness = [
        f"MVP Ready: {overall.get('is_mvp_ready','—')}",
        f"Investment Ready: {overall.get('is_investment_ready','—')}",
        f"Incubator Ready: {overall.get('is_incubator_ready','—')}",
    ]
    _bullet_box(slide, readiness,
                Inches(0.4), Inches(0.85), W - Inches(0.8), Inches(0.9),
                bullet="", color=ACCENT2, font_size=13)

    _add_textbox(slide, "Team & Funding",
                 Inches(0.4), Inches(1.85), Inches(3.5), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _add_textbox(slide,
                 f"{support.get('team_needed','—')}\n{support.get('funding_stage','—')}",
                 Inches(0.4), Inches(2.2), Inches(3.5), Inches(0.75),
                 font_size=12, color=BODY_COL)

    _add_textbox(slide, "Partnerships",
                 Inches(4.2), Inches(1.85), Inches(3.4), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    _add_textbox(slide, support.get("partnerships", "—"),
                 Inches(4.2), Inches(2.2), Inches(3.4), Inches(0.75),
                 font_size=12, color=BODY_COL)

    _add_textbox(slide, "Recommended Tech Stack",
                 Inches(0.4), Inches(3.1), W - Inches(0.8), Inches(0.35),
                 font_size=13, bold=True, color=ACCENT)
    stack = [
        f"Backend: {tech.get('backend','—')}",
        f"Frontend: {tech.get('frontend','—')}",
        f"Database: {tech.get('database','—')}",
        f"AI Tools: {tech.get('ai_tools','—')}",
    ]
    _bullet_box(slide, stack,
                Inches(0.4), Inches(3.45), W - Inches(0.8), Inches(1.4),
                bullet="▸", font_size=12)


# ─── Main entry point (same signature as original) ───────────────────────────

def generate_ppt(slide_data: dict) -> bytes:
    """
    Generates a complete pitch deck from the analysis dict.
    No template file required — builds all slides programmatically.

    Accepts either:
    - The full analysis dict (from analyze_idea)
    - A slides dict from generate_pitch_slides (treated as analysis fallback)
    """
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.63)   # 16:9

    _slide_title(prs, slide_data)
    _slide_problem(prs, slide_data)
    _slide_solution(prs, slide_data)
    _slide_market(prs, slide_data)
    _slide_scores(prs, slide_data)
    _slide_roadmap(prs, slide_data)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()